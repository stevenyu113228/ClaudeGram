"""File handling for Telegram file uploads."""
import base64
import hashlib
import logging
import os
from dataclasses import dataclass
from enum import Enum
from typing import Optional

import aiohttp

logger = logging.getLogger(__name__)


class FileType(Enum):
    """Supported file types."""

    IMAGE = "image"
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    UNSUPPORTED = "unsupported"


@dataclass
class ProcessedFile:
    """Processed file ready for Claude API."""

    file_type: FileType
    file_name: str
    mime_type: str
    file_size: int
    telegram_file_id: str
    content_hash: str
    # For images/PDFs - base64 encoded data
    base64_data: Optional[str] = None
    # For DOCX/PPTX - extracted text
    extracted_text: Optional[str] = None


# MIME type mappings
MIME_TO_FILE_TYPE = {
    "image/jpeg": FileType.IMAGE,
    "image/png": FileType.IMAGE,
    "image/gif": FileType.IMAGE,
    "image/webp": FileType.IMAGE,
    "application/pdf": FileType.PDF,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": FileType.DOCX,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": FileType.PPTX,
}

SUPPORTED_EXTENSIONS = {
    ".jpg": FileType.IMAGE,
    ".jpeg": FileType.IMAGE,
    ".png": FileType.IMAGE,
    ".gif": FileType.IMAGE,
    ".webp": FileType.IMAGE,
    ".pdf": FileType.PDF,
    ".docx": FileType.DOCX,
    ".pptx": FileType.PPTX,
}

# File size limits (in bytes)
MAX_FILE_SIZE = 20 * 1024 * 1024  # 20MB (Telegram Bot API limit)


def detect_file_type(file_name: str | None, mime_type: str | None) -> FileType:
    """
    Detect file type from MIME type or extension.

    Args:
        file_name: Original file name
        mime_type: MIME type if known

    Returns:
        FileType enum value
    """
    # Try MIME type first
    if mime_type and mime_type in MIME_TO_FILE_TYPE:
        return MIME_TO_FILE_TYPE[mime_type]

    # Fall back to extension
    if file_name:
        ext = os.path.splitext(file_name.lower())[1]
        if ext in SUPPORTED_EXTENSIONS:
            return SUPPORTED_EXTENSIONS[ext]

    return FileType.UNSUPPORTED


async def download_telegram_file(
    bot_token: str,
    file_id: str,
) -> tuple[bytes, str]:
    """
    Download file from Telegram servers.

    Args:
        bot_token: Telegram bot token
        file_id: Telegram file ID

    Returns:
        Tuple of (file_bytes, file_path)

    Raises:
        ValueError: If file download fails
    """
    async with aiohttp.ClientSession() as session:
        # Get file path from Telegram
        get_file_url = f"https://api.telegram.org/bot{bot_token}/getFile"
        async with session.get(
            get_file_url,
            params={"file_id": file_id},
            timeout=aiohttp.ClientTimeout(total=30),
        ) as resp:
            result = await resp.json()
            if not result.get("ok"):
                raise ValueError(f"Failed to get file info: {result}")
            file_path = result["result"]["file_path"]
            file_size = result["result"].get("file_size", 0)

        # Check file size
        if file_size > MAX_FILE_SIZE:
            raise ValueError(
                f"File too large: {file_size / 1024 / 1024:.1f}MB (max {MAX_FILE_SIZE / 1024 / 1024:.0f}MB)"
            )

        # Download file content
        download_url = f"https://api.telegram.org/file/bot{bot_token}/{file_path}"
        async with session.get(
            download_url,
            timeout=aiohttp.ClientTimeout(total=60),
        ) as resp:
            if resp.status != 200:
                raise ValueError(f"Failed to download file: HTTP {resp.status}")
            file_bytes = await resp.read()

    return file_bytes, file_path


def extract_text_from_docx(file_bytes: bytes) -> str:
    """
    Extract text from DOCX file.

    Args:
        file_bytes: Raw file bytes

    Returns:
        Extracted text content
    """
    from docx import Document
    import io

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = []

    # Extract paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                paragraphs.append(" | ".join(row_text))

    return "\n\n".join(paragraphs)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract text from PPTX file.

    Args:
        file_bytes: Raw file bytes

    Returns:
        Extracted text content
    """
    from pptx import Presentation
    import io

    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if len(slide_content) > 1:  # Only add if there's content beyond the header
            slides_text.append("\n".join(slide_content))

    return "\n\n".join(slides_text)


async def process_file(
    bot_token: str,
    file_id: str,
    file_name: str | None,
    mime_type: str | None,
    file_size: int,
) -> ProcessedFile:
    """
    Download and process a file from Telegram.

    Args:
        bot_token: Telegram bot token
        file_id: Telegram file ID
        file_name: Original file name
        mime_type: MIME type if known
        file_size: File size in bytes

    Returns:
        ProcessedFile with base64 data or extracted text

    Raises:
        ValueError: If file type is unsupported or processing fails
    """
    file_type = detect_file_type(file_name, mime_type)

    if file_type == FileType.UNSUPPORTED:
        raise ValueError(f"Unsupported file type: {file_name} ({mime_type})")

    # Download file
    logger.info(f"Downloading file: {file_name} ({file_type.value})")
    file_bytes, file_path = await download_telegram_file(bot_token, file_id)

    # Generate content hash
    content_hash = hashlib.md5(file_bytes).hexdigest()

    # Infer mime_type from file_path if not provided
    if not mime_type:
        ext = os.path.splitext(file_path.lower())[1]
        mime_map = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }
        mime_type = mime_map.get(ext, "application/octet-stream")

    # Process based on file type
    if file_type in (FileType.IMAGE, FileType.PDF):
        # Native Claude support - encode as base64
        logger.info(f"Encoding {file_type.value} as base64")
        base64_data = base64.standard_b64encode(file_bytes).decode("utf-8")
        return ProcessedFile(
            file_type=file_type,
            file_name=file_name or os.path.basename(file_path),
            mime_type=mime_type,
            file_size=len(file_bytes),
            telegram_file_id=file_id,
            content_hash=content_hash,
            base64_data=base64_data,
        )

    elif file_type == FileType.DOCX:
        logger.info("Extracting text from DOCX")
        try:
            extracted_text = extract_text_from_docx(file_bytes)
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise ValueError(f"Failed to extract text from DOCX: {e}")
        return ProcessedFile(
            file_type=file_type,
            file_name=file_name or "document.docx",
            mime_type=mime_type,
            file_size=len(file_bytes),
            telegram_file_id=file_id,
            content_hash=content_hash,
            extracted_text=extracted_text,
        )

    elif file_type == FileType.PPTX:
        logger.info("Extracting text from PPTX")
        try:
            extracted_text = extract_text_from_pptx(file_bytes)
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise ValueError(f"Failed to extract text from PPTX: {e}")
        return ProcessedFile(
            file_type=file_type,
            file_name=file_name or "presentation.pptx",
            mime_type=mime_type,
            file_size=len(file_bytes),
            telegram_file_id=file_id,
            content_hash=content_hash,
            extracted_text=extracted_text,
        )

    raise ValueError(f"Unhandled file type: {file_type}")


def get_supported_formats_message() -> str:
    """Get a user-friendly message about supported file formats."""
    return "支援的檔案格式：JPG、PNG、GIF、WEBP（圖片）、PDF、DOCX（Word）、PPTX（PowerPoint）"
